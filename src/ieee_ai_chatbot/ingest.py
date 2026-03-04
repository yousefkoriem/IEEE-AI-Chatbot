from __future__ import annotations

import hashlib
import importlib
import json
from pathlib import Path
from collections import deque
from typing import Any
from urllib.parse import urljoin, urlparse

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader
from pptx import Presentation
import requests
from bs4 import BeautifulSoup

from .config import Settings
from .vectorstore import get_vector_store

SUPPORTED_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".docx", ".doc"}


def _sha256_file(path: Path) -> str:
    hasher = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(65536), b""):
            hasher.update(block)
    return hasher.hexdigest()


def _load_manifest(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {"sources": {}}
    return json.loads(path.read_text(encoding="utf-8"))


def _save_manifest(path: Path, payload: dict[str, Any]) -> None:
    path.write_text(json.dumps(payload, indent=2), encoding="utf-8")


def _extract_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    pages: list[str] = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n".join(pages).strip()


def _extract_ppt_text(path: Path) -> str:
    presentation = Presentation(str(path))
    blocks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                blocks.append(text)
    return "\n".join(blocks).strip()


def _extract_docx_text(path: Path) -> str:
    docx_module = importlib.import_module("docx")
    document = docx_module.Document(str(path))
    blocks: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            blocks.append(text)
    return "\n".join(blocks).strip()


def _extract_doc_text(path: Path) -> str:
    try:
        from unstructured.partition.auto import partition
    except Exception as error:
        raise RuntimeError(
            "DOC parsing requires unstructured document support."
        ) from error

    elements = partition(filename=str(path))
    blocks = [str(element).strip() for element in elements if str(element).strip()]
    return "\n".join(blocks).strip()


def _extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf_text(path)
    if suffix in {".ppt", ".pptx"}:
        return _extract_ppt_text(path)
    if suffix == ".docx":
        return _extract_docx_text(path)
    if suffix == ".doc":
        return _extract_doc_text(path)
    raise ValueError(f"Unsupported file extension: {path.suffix}")


def _build_documents(path: Path, source_id: str, text: str) -> list[Document]:
    return [
        Document(
            page_content=text,
            metadata={
                "source": source_id,
                "filename": path.name,
                "suffix": path.suffix.lower(),
            },
        )
    ]


def _website_chunks_to_documents(url: str, text: str, title: str) -> list[Document]:
    return [
        Document(
            page_content=text,
            metadata={
                "source": url,
                "url": url,
                "title": title,
                "suffix": ".html",
                "filename": url,
            },
        )
    ]


def _extract_page_text(html: str) -> tuple[str, str]:
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title = (soup.title.string or "").strip() if soup.title else ""
    text = soup.get_text("\n", strip=True)
    return text, title


def _crawl_same_domain(start_url: str, max_pages: int, timeout_seconds: int) -> list[dict[str, str]]:
    parsed_start = urlparse(start_url)
    base_domain = parsed_start.netloc
    if not parsed_start.scheme:
        start_url = f"https://{start_url}"

    queue: deque[str] = deque([start_url])
    visited: set[str] = set()
    pages: list[dict[str, str]] = []

    headers = {
        "User-Agent": "IEEE-AI-Chatbot-RAG/1.0 (+https://ieee-mangment.vercel.app/)"
    }

    while queue and len(pages) < max_pages:
        url = queue.popleft()
        if url in visited:
            continue
        visited.add(url)

        try:
            response = requests.get(url, timeout=timeout_seconds, headers=headers)
            if response.status_code >= 400:
                continue
            if "text/html" not in response.headers.get("Content-Type", ""):
                continue
        except requests.RequestException:
            continue

        text, title = _extract_page_text(response.text)
        if text:
            pages.append({"url": url, "text": text, "title": title})

        soup = BeautifulSoup(response.text, "html.parser")
        for a_tag in soup.find_all("a", href=True):
            absolute = urljoin(url, a_tag["href"]).split("#", 1)[0]
            parsed = urlparse(absolute)
            if parsed.scheme not in {"http", "https"}:
                continue
            if parsed.netloc != base_domain:
                continue
            if absolute not in visited:
                queue.append(absolute)

    return pages


def ingest_files(settings: Settings, file_paths: list[str], origin: str) -> dict[str, int]:
    vector_store = get_vector_store(settings)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
    )
    manifest_path = Path(settings.manifest_path)
    manifest = _load_manifest(manifest_path)
    sources = manifest.setdefault("sources", {})

    indexed = 0
    skipped = 0
    deleted = 0

    for raw_path in file_paths:
        path = Path(raw_path)
        if not path.exists() or path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            skipped += 1
            continue

        source_id = str(path.resolve())
        file_hash = _sha256_file(path)
        existing = sources.get(source_id)
        if existing and existing.get("hash") == file_hash:
            skipped += 1
            continue

        text = _extract_text(path)
        if not text:
            skipped += 1
            continue

        documents = _build_documents(path, source_id, text)
        chunks = splitter.split_documents(documents)
        chunk_ids: list[str] = []
        for idx, chunk in enumerate(chunks):
            chunk_id = hashlib.sha1(f"{source_id}:{file_hash}:{idx}".encode("utf-8")).hexdigest()
            chunk.metadata["chunk_id"] = chunk_id
            chunk.metadata["origin"] = origin
            chunk.metadata["hash"] = file_hash
            chunk_ids.append(chunk_id)

        if existing and existing.get("chunk_ids"):
            vector_store.delete(ids=existing["chunk_ids"])
            deleted += len(existing["chunk_ids"])

        vector_store.add_documents(documents=chunks, ids=chunk_ids)
        sources[source_id] = {
            "hash": file_hash,
            "chunk_ids": chunk_ids,
            "origin": origin,
        }
        indexed += 1

    _save_manifest(manifest_path, manifest)
    return {"indexed": indexed, "skipped": skipped, "deleted": deleted}


def sync_local_docs(settings: Settings) -> dict[str, int]:
    search_roots = [
        Path(settings.docs_pdf_dir),
        Path(settings.docs_ppt_dir),
        Path(settings.docs_doc_dir),
    ]
    all_files: list[str] = []
    for root in search_roots:
        if not root.exists():
            continue
        for file in root.rglob("*"):
            if file.is_file() and file.suffix.lower() in SUPPORTED_EXTENSIONS:
                all_files.append(str(file))

    result = ingest_files(settings, all_files, origin="local")

    manifest_path = Path(settings.manifest_path)
    manifest = _load_manifest(manifest_path)
    sources = manifest.setdefault("sources", {})

    active_sources = {str(Path(path).resolve()) for path in all_files}
    vector_store = get_vector_store(settings)
    delete_count = 0
    removed_sources = []

    for source_id, meta in list(sources.items()):
        if meta.get("origin") != "local":
            continue
        if source_id in active_sources:
            continue
        chunk_ids = meta.get("chunk_ids", [])
        if chunk_ids:
            vector_store.delete(ids=chunk_ids)
            delete_count += len(chunk_ids)
        removed_sources.append(source_id)

    for source_id in removed_sources:
        sources.pop(source_id, None)

    _save_manifest(manifest_path, manifest)
    result["deleted"] += delete_count
    result["total_files"] = len(all_files)
    return result


def ingest_website(settings: Settings, start_url: str, max_pages: int = 25) -> dict[str, int]:
    pages = _crawl_same_domain(
        start_url=start_url,
        max_pages=max_pages,
        timeout_seconds=settings.website_timeout_seconds,
    )

    vector_store = get_vector_store(settings)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
    )
    manifest_path = Path(settings.manifest_path)
    manifest = _load_manifest(manifest_path)
    sources = manifest.setdefault("sources", {})

    indexed = 0
    skipped = 0
    deleted = 0

    for page in pages:
        url = page["url"]
        text = page["text"]
        title = page["title"]
        content_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()

        existing = sources.get(url)
        if existing and existing.get("hash") == content_hash:
            skipped += 1
            continue

        docs = _website_chunks_to_documents(url=url, text=text, title=title)
        chunks = splitter.split_documents(docs)
        chunk_ids: list[str] = []
        for idx, chunk in enumerate(chunks):
            chunk_id = hashlib.sha1(f"{url}:{content_hash}:{idx}".encode("utf-8")).hexdigest()
            chunk.metadata["chunk_id"] = chunk_id
            chunk.metadata["origin"] = "website"
            chunk.metadata["hash"] = content_hash
            chunk_ids.append(chunk_id)

        if existing and existing.get("chunk_ids"):
            vector_store.delete(ids=existing["chunk_ids"])
            deleted += len(existing["chunk_ids"])

        vector_store.add_documents(documents=chunks, ids=chunk_ids)
        sources[url] = {
            "hash": content_hash,
            "chunk_ids": chunk_ids,
            "origin": "website",
        }
        indexed += 1

    parsed_start = urlparse(start_url if "://" in start_url else f"https://{start_url}")
    crawled_urls = {page["url"] for page in pages}
    stale_sources: list[str] = []

    for source_id, meta in list(sources.items()):
        if meta.get("origin") != "website":
            continue

        parsed_source = urlparse(source_id)
        if parsed_source.netloc != parsed_start.netloc:
            continue
        if source_id in crawled_urls:
            continue

        chunk_ids = meta.get("chunk_ids", [])
        if chunk_ids:
            vector_store.delete(ids=chunk_ids)
            deleted += len(chunk_ids)
        stale_sources.append(source_id)

    for source_id in stale_sources:
        sources.pop(source_id, None)

    _save_manifest(manifest_path, manifest)
    return {
        "indexed": indexed,
        "skipped": skipped,
        "deleted": deleted,
        "total_pages": len(pages),
    }
