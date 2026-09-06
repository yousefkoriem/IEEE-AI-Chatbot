"""Document ingestion and indexing entry points."""

import logging
from pathlib import Path
from typing import Any

from langchain_text_splitters import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# File loaders – each returns a list of dicts with 'text' and metadata keys
# ---------------------------------------------------------------------------

def _load_text_file(path: Path) -> list[dict[str, Any]]:
    """Load plain text, markdown, or HTML as raw text."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return [{"text": text, "source": path.name}]


def _load_pdf(path: Path) -> list[dict[str, Any]]:
    """Load a PDF file page-by-page using PyPDF."""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[dict[str, Any]] = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append({"text": text, "source": path.name, "page": i + 1})
    return pages


def _load_docx(path: Path) -> list[dict[str, Any]]:
    """Load a DOCX file paragraph-by-paragraph."""
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"text": full_text, "source": path.name}] if full_text else []


def _load_pptx(path: Path) -> list[dict[str, Any]]:
    """Load a PPTX file slide-by-slide."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slide_text = "\n".join(texts)
        if slide_text.strip():
            slides.append({"text": slide_text, "source": path.name, "page": i + 1})
    return slides


# Map of extensions to loader functions
_LOADERS = {
    ".md": _load_text_file,
    ".txt": _load_text_file,
    ".html": _load_text_file,
    ".htm": _load_text_file,
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
}

SUPPORTED_EXTENSIONS = set(_LOADERS.keys())


# ---------------------------------------------------------------------------
# Main ingestion pipeline
# ---------------------------------------------------------------------------

def ingest_documents(
    source_dir: str | Path,
    vectorstore=None,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> int:
    """Ingest documents from *source_dir* into the vector store.

    Supports: .md, .txt, .html, .pdf, .docx, .pptx

    Returns the number of chunks indexed.
    """
    source_path = Path(source_dir)
    if not source_path.exists():
        raise FileNotFoundError(f"Source directory not found: {source_path}")

    # Discover files
    files = [
        f
        for f in source_path.rglob("*")
        if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS
    ]
    if not files:
        logger.warning("No supported files found in %s", source_path)
        return 0

    logger.info("Found %d files to ingest in %s", len(files), source_path)

    # Load all documents
    raw_docs: list[dict[str, Any]] = []
    for file_path in files:
        loader = _LOADERS.get(file_path.suffix.lower())
        if loader is None:
            continue
        try:
            docs = loader(file_path)
            # Enrich metadata
            category = _extract_category(file_path, source_path)
            title = _extract_title(file_path)
            for doc in docs:
                doc["category"] = category
                doc["title"] = title
            raw_docs.extend(docs)
            logger.info("Loaded %s (%d segments)", file_path.name, len(docs))
        except Exception:
            logger.exception("Failed to load %s", file_path)

    if not raw_docs:
        logger.warning("No text extracted from any files")
        return 0

    # Chunk documents
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    chunks: list[dict[str, Any]] = []
    for doc in raw_docs:
        split_texts = splitter.split_text(doc["text"])
        for i, chunk_text in enumerate(split_texts):
            chunks.append(
                {
                    "text": chunk_text,
                    "source": doc.get("source", ""),
                    "category": doc.get("category", ""),
                    "title": doc.get("title", ""),
                    "page": doc.get("page", 0),
                    "chunk_index": i,
                }
            )

    logger.info("Split %d raw segments into %d chunks", len(raw_docs), len(chunks))

    # Upsert to vector store
    if vectorstore is not None:
        count = vectorstore.add_documents(chunks)
        logger.info("Indexed %d chunks into Pinecone", count)
        return count

    logger.info("No vectorstore provided – returning chunk count only")
    return len(chunks)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _extract_category(file_path: Path, base_path: Path) -> str:
    """Derive a category from the file's parent directory relative to base."""
    try:
        relative = file_path.relative_to(base_path)
        parts = relative.parts[:-1]  # exclude filename
        return "/".join(parts) if parts else "general"
    except ValueError:
        return "general"


def _extract_title(file_path: Path) -> str:
    """Extract a title from the first heading or use the filename."""
    if file_path.suffix.lower() in {".md", ".txt"}:
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                for line in f:
                    line = line.strip()
                    if line.startswith("# "):
                        return line.lstrip("# ").strip()
                    if line:
                        return line[:80]
        except Exception:
            pass
    return file_path.stem.replace("_", " ").replace("-", " ").title()
